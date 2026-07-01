# -*- coding: utf-8 -*-
"""Génération de PowerPoint à partir des données KPI."""
# ── RECONSTRUIT ──


def generate_pptx(kpi_data, output_path="dashboard_kpi.pptx"):
    """Génère un PowerPoint résumant les KPI."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Slide titre
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # Fond bleu
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "DASHBOARD KPI"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        date_str = kpi_data.get("date", "")
        if date_str:
            p2 = tf.add_paragraph()
            p2.text = date_str
            p2.font.size = Pt(24)
            p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
            p2.alignment = PP_ALIGN.CENTER

        # Slide Performance
        if kpi_data.get("p_rows"):
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            _add_kpi_table(slide2, "PERFORMANCE", kpi_data["p_rows"], kpi_data.get("p_cols", []))

        # Slide Qualité
        if kpi_data.get("q_rows"):
            slide3 = prs.slides.add_slide(prs.slide_layouts[6])
            _add_kpi_table(slide3, "QUALITE", kpi_data["q_rows"], kpi_data.get("q_cols", []))

        # Slide Anomalies
        ano_rows = kpi_data.get("ano_p_r", []) + kpi_data.get("ano_q_r", [])
        if ano_rows:
            slide4 = prs.slides.add_slide(prs.slide_layouts[6])
            _add_kpi_table(slide4, "ANOMALIES", ano_rows,
                           kpi_data.get("ano_p_c", ["Poste", "KPI", "Valeur", "Cible", "Statut"]))

        prs.save(output_path)
        return output_path
    except ImportError:
        return None


def _add_kpi_table(slide, title, rows, cols):
    """Ajoute un tableau KPI à une slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Titre
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

    if not cols or not rows:
        return

    n_cols = len(cols)
    n_rows = min(len(rows), 20)  # Limiter à 20 lignes

    table_shape = slide.shapes.add_table(
        n_rows + 1, n_cols,
        Inches(0.3), Inches(1.2),
        Inches(12.7), Inches(0.3 * (n_rows + 1))
    )
    table = table_shape.table

    # En-têtes
    for j, col_name in enumerate(cols):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True

    # Données
    for i, row in enumerate(rows[:n_rows]):
        for j, col_name in enumerate(cols):
            cell = table.cell(i + 1, j)
            val = row.get(col_name, "")
            cell.text = str(val) if val is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
