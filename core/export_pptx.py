# -*- coding: utf-8 -*-
"""
Génération d'une présentation PowerPoint dynamique selon le filtre
Poste de travail (SF1 → Maroc Chimie, SF2 → FEEDS, mixte → OCP).

Usage dans app.py :
    from core.export_pptx import build_presentation
    pptx_bytes = build_presentation(vp, ckdf, ano_map, pa, qa,
                                    pscores, qscores, hist_df, fichier_date)
    st.download_button("📊 Exporter PowerPoint", pptx_bytes,
                       file_name="presentation_kpis.pptx")
"""
import io
from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from core.constants import QK, PK, CIBLE, KPI_RESP_MAP

# ── Palette OCP ──────────────────────────────────────────────────────────
GREEN_OCP  = RGBColor(0x2C, 0x5F, 0x2D)
MOSS       = RGBColor(0x97, 0xBC, 0x62)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1E, 0x29, 0x3B)
GREY       = RGBColor(0x64, 0x74, 0x8B)
RED        = RGBColor(0xEF, 0x44, 0x44)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
LIGHT_BG   = RGBColor(0xF5, 0xF5, 0xF5)

SW, SH = Inches(13.333), Inches(7.5)   # 16:9


def _entity_name(vp):
    """SF1 → Maroc Chimie, SF2 → FEEDS, mixte → OCP (SF1 & SF2)."""
    has_sf1 = any(str(p).startswith("SF1") for p in vp)
    has_sf2 = any(str(p).startswith("SF2") for p in vp)
    if has_sf1 and not has_sf2:
        return "Maroc Chimie"
    if has_sf2 and not has_sf1:
        return "FEEDS"
    return "OCP — Maroc Chimie & FEEDS"


def _score_color(v, s1=70, s2=90):
    if v >= s2: return GREEN
    if v >= s1: return AMBER
    return RED


def _add_title_bar(slide, text, subtitle=None):
    """Bandeau titre haut de slide (fond vert OCP)."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = GREEN_OCP
    r.font.name = "Calibri"
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(14); r2.font.color.rgb = GREY; r2.font.name = "Calibri"


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ═══════════════════════════════════════════════════════════════════════════
def _slide_title(prs, entity, fichier_date):
    slide = _blank_slide(prs)
    # Fond vert plein
    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = GREEN_OCP
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Bande moss décorative en bas
    band = slide.shapes.add_shape(1, 0, Inches(6.6), SW, Inches(0.9))
    band.fill.solid(); band.fill.fore_color.rgb = MOSS
    band.line.fill.background(); band.shadow.inherit = False

    # Titre principal
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.3), Inches(2.2))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Présentation des KPIs"
    r.font.size = Pt(48); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "Maintenance SAP PM"
    r2.font.size = Pt(30); r2.font.color.rgb = RGBColor(0xE7, 0xE8, 0xD1); r2.font.name = "Calibri"

    # Entité
    box2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.6), Inches(11.3), Inches(1.0))
    tf2 = box2.text_frame; p3 = tf2.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = entity
    r3.font.size = Pt(32); r3.font.bold = True; r3.font.color.rgb = WHITE; r3.font.name = "Calibri"

    # Date
    box3 = slide.shapes.add_textbox(Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.7))
    tf3 = box3.text_frame; p4 = tf3.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run(); r4.text = f"Données du {fichier_date}   •   Généré le {datetime.now().strftime('%d/%m/%Y')}"
    r4.font.size = Pt(13); r4.font.color.rgb = DARK; r4.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════════
# Table helper
# ═══════════════════════════════════════════════════════════════════════════
def _add_table(slide, headers, rows, left, top, width, height,
               col_widths=None, header_fill=GREEN_OCP, font_size=10,
               color_col=None, cible_map=None, lower_set=None):
    nrows, ncols = len(rows) + 1, len(headers)
    tbl_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

    # Header
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = header_fill
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = c.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(h)
        r.font.size = Pt(font_size); r.font.bold = True; r.font.color.rgb = WHITE
        r.font.name = "Calibri"

    # Body
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT_BG
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(font_size); r.font.name = "Calibri"; r.font.color.rgb = DARK
            # Colorer la colonne des valeurs selon score/cible
            if color_col is not None and j == color_col:
                try:
                    fv = float(str(val).replace('%', '').replace(',', '.'))
                    kpi_name = str(row[0])
                    tgt = (cible_map or {}).get(kpi_name, 100)
                    lower = kpi_name in (lower_set or set())
                    ok = (fv <= tgt) if lower else (fv >= tgt)
                    r.font.color.rgb = GREEN if ok else RED
                    r.font.bold = True
                except (ValueError, TypeError):
                    pass
    return tbl_shape


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Scores globaux + Indicateurs Performance (barres)
# ═══════════════════════════════════════════════════════════════════════════
def _bar_row(slide, label, value, y, x0=Inches(0.8), maxw=8.5,
             s1=70, s2=90, label_w=3.2):
    """Une barre horizontale + label + valeur."""
    # Label
    lb = slide.shapes.add_textbox(x0, y, Inches(label_w), Inches(0.32))
    tf = lb.text_frame; tf.margin_left = 0; tf.margin_top = 0
    p = tf.paragraphs[0]; r = p.add_run(); r.text = str(label)
    r.font.size = Pt(11); r.font.color.rgb = DARK; r.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT

    # Barre fond
    bar_x = x0 + Inches(label_w + 0.15)
    track = slide.shapes.add_shape(1, bar_x, y + Emu(20000),
                                    Inches(maxw), Inches(0.20))
    track.fill.solid(); track.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    track.line.fill.background(); track.shadow.inherit = False

    # Barre valeur
    w = max(0.02, maxw * min(value, 100) / 100.0)
    bar = slide.shapes.add_shape(1, bar_x, y + Emu(20000),
                                  Inches(w), Inches(0.20))
    bar.fill.solid(); bar.fill.fore_color.rgb = _score_color(value, s1, s2)
    bar.line.fill.background(); bar.shadow.inherit = False

    # Valeur
    vb = slide.shapes.add_textbox(bar_x + Inches(maxw + 0.1), y, Inches(0.9), Inches(0.32))
    tf2 = vb.text_frame; tf2.margin_left = 0; tf2.margin_top = 0
    p2 = tf2.paragraphs[0]; r2 = p2.add_run(); r2.text = f"{value:.0f}%"
    r2.font.size = Pt(11); r2.font.bold = True; r2.font.color.rgb = DARK; r2.font.name = "Calibri"


def _slide_scores_and_perf(prs, entity, vp, pscores, qscores, pa):
    slide = _blank_slide(prs)
    _add_title_bar(slide, "Scores globaux par poste",
                   f"{entity} — Performance (vert) et Qualité (bleu)")

    # Section gauche : scores par poste
    postes = [p for p in vp if p in pscores or p in qscores][:9]
    y = Inches(1.5)
    for poste in postes:
        ps = pscores.get(poste, 0)
        _bar_row(slide, poste, ps, y, x0=Inches(0.4), maxw=4.2, label_w=1.9)
        y += Inches(0.55)

    # Section droite : Taux moyens Performance
    hdr = slide.shapes.add_textbox(Inches(7.2), Inches(1.35), Inches(5.5), Inches(0.4))
    p = hdr.text_frame.paragraphs[0]; r = p.add_run()
    r.text = "Taux moyens — Performance"
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = GREEN_OCP; r.font.name = "Calibri"

    y2 = Inches(1.9)
    for kpi in QK[:9]:
        if kpi in pa:
            short = kpi.replace("OT ", "").replace("Performance ", "Perf ")[:22]
            _bar_row(slide, short, pa[kpi], y2, x0=Inches(6.7), maxw=3.2, label_w=2.4)
            y2 += Inches(0.52)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDES 3 & 4 — Détail indicateurs + anomalies (Perf / Qualité)
# ═══════════════════════════════════════════════════════════════════════════
def _slide_detail(prs, entity, vp, ckdf, ano_map, kpi_list, kind):
    slide = _blank_slide(prs)
    color = GREEN_OCP
    _add_title_bar(slide, f"Détail des indicateurs de {kind}",
                   f"{entity} — valeurs et nombre d'anomalies par KPI")

    postes = [p for p in vp if p in ckdf.index][:8]

    # Tableau 1 : valeurs KPI par poste (moyenne)
    headers = ["Indicateur", "Valeur moy.", "Cible", "Anomalies"]
    rows = []
    for kpi in kpi_list:
        # Valeur moyenne sur les postes
        vals = [float(ckdf.loc[p, kpi]) for p in postes if kpi in ckdf.columns and p in ckdf.index]
        vmoy = sum(vals) / len(vals) if vals else 0
        tgt = CIBLE.get(kpi, 100)
        # Total anomalies
        nb = 0
        if kpi in ano_map:
            s = ano_map[kpi]
            nb = int(sum(int(s.get(p, 0)) for p in postes))
        short = kpi.replace("TAUX_REALISATION_CORRECTIF/PT", "Taux Réalis. Correctif")
        rows.append([short[:38], f"{vmoy:.1f}%", f"{tgt}%", str(nb)])

    _add_table(
        slide, headers, rows,
        Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.3),
        col_widths=[Inches(6.5), Inches(2.0), Inches(1.8), Inches(2.0)],
        header_fill=color, font_size=11,
        color_col=1, cible_map=CIBLE,
        lower_set=set(k for k in kpi_list if k in _LOWER),
    )


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Sparklines (évolution par poste, version texte tendance)
# ═══════════════════════════════════════════════════════════════════════════
def _slide_sparklines(prs, entity, vp, hist_df, pscores, qscores):
    slide = _blank_slide(prs)
    _add_title_bar(slide, "Suivi Sparklines par Poste de Travail",
                   f"{entity} — évolution récente des scores")

    postes = [p for p in vp if p in pscores or p in qscores][:10]
    headers = ["Poste de travail", "Performance", "Qualité", "Tendance"]
    rows = []
    for poste in postes:
        ps = pscores.get(poste, 0)
        qs = qscores.get(poste, 0)
        # Tendance simple depuis hist_df si dispo
        trend = "→"
        if hist_df is not None and not hist_df.empty:
            try:
                h = hist_df[hist_df["Poste"] == poste].sort_values("Date")
                if len(h) >= 2:
                    d = h.iloc[-1]["Value"] - h.iloc[-2]["Value"]
                    trend = "↑" if d > 1 else ("↓" if d < -1 else "→")
            except Exception:
                pass
        rows.append([poste, f"{ps:.0f}%", f"{qs:.0f}%", trend])

    _add_table(
        slide, headers, rows,
        Inches(1.2), Inches(1.6), Inches(10.9), Inches(5.0),
        col_widths=[Inches(4.5), Inches(2.3), Inches(2.3), Inches(1.8)],
        header_fill=GREEN_OCP, font_size=12,
    )


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Plan d'action
# ═══════════════════════════════════════════════════════════════════════════
def _slide_plan_action(prs, entity, vp, ckdf, ano_map):
    slide = _blank_slide(prs)
    _add_title_bar(slide, "Plan d'action",
                   f"{entity} — KPIs nécessitant une action (anomalies > 0)")

    postes = [p for p in vp if p in ckdf.index]
    rows = []
    for poste in postes:
        for kpi in list(QK) + list(PK):
            if kpi not in ano_map:
                continue
            nb = int(ano_map[kpi].get(poste, 0))
            if nb > 0:
                resp = KPI_RESP_MAP.get(kpi, "Non assigné")
                short = kpi.replace("TAUX_REALISATION_CORRECTIF/PT", "Taux Réalis. Corr.")
                rows.append([poste, short[:30], str(nb), resp[:20]])
    # Trier par nb anomalies décroissant, limiter à 12 lignes
    rows.sort(key=lambda r: -int(r[2]))
    rows = rows[:12]

    if not rows:
        box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
        p = box.text_frame.paragraphs[0]; r = p.add_run()
        r.text = "🎉 Aucune action requise — tous les indicateurs sont conformes."
        r.font.size = Pt(20); r.font.color.rgb = GREEN; r.font.name = "Calibri"
        return

    headers = ["Poste", "Indicateur", "Anomalies", "Responsable"]
    _add_table(
        slide, headers, rows,
        Inches(0.6), Inches(1.55), Inches(12.1), Inches(5.4),
        col_widths=[Inches(2.4), Inches(5.0), Inches(1.9), Inches(2.8)],
        header_fill=RED, font_size=11,
    )


# LOWER_BETTER (import paresseux pour éviter les erreurs si absent)
try:
    from core.constants import LOWER_BETTER as _LOWER
except Exception:
    _LOWER = []


# ═══════════════════════════════════════════════════════════════════════════
# API publique
# ═══════════════════════════════════════════════════════════════════════════
def build_presentation(vp, ckdf, ano_map, pa, qa,
                        pscores, qscores, hist_df=None, fichier_date=""):
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    entity = _entity_name(vp)

    _slide_title(prs, entity, fichier_date)
    _slide_scores_and_perf(prs, entity, vp, pscores, qscores, pa)
    _slide_detail(prs, entity, vp, ckdf, ano_map, QK, "Performance")
    _slide_detail(prs, entity, vp, ckdf, ano_map, PK, "Qualité")
    _slide_sparklines(prs, entity, vp, hist_df, pscores, qscores)
    _slide_plan_action(prs, entity, vp, ckdf, ano_map)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
