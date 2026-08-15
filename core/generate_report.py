# -*- coding: utf-8 -*-
"""
Génère un rapport KPI (1 slide) par poste de travail, en Python pur
(python-pptx) — pas de dépendance Node.js, exécutable directement dans
l'app Streamlit après chaque recalcul.

Reproduit le design déjà validé (voir gen_sf1_ecu_1page.js) :
- En-tête : poste, date, 3 badges (Score Performance / Score Qualité / Anomalies)
- 2 tableaux KPI compacts (Performance / Qualité), sans colonne "Anomalies"
- Graphique (barres horizontales bleues) : anomalies par indicateur, noms abrégés
- Bloc "Suivi d'évolution" (1 date pour l'instant)
- Tableau "Plan d'action" : TOUS les indicateurs en anomalie, avec colonne
  "Nécessité action" (Oui=rouge / Non=vert) à la place du nombre d'anomalies
- Crédit applicatif sous l'en-tête
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

from core.constants import LOWER_BETTER

# ── Palette ──
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GREY = RGBColor(0x64, 0x74, 0x8B)
LGREY = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GREEN_ACCENT = RGBColor(0x05, 0x96, 0x69)
BLUE_ACCENT = RGBColor(0x25, 0x63, 0xEB)

SHORT_LABELS = {
    "TAUX_REALISATION_CORRECTIF/PT": "Taux Réal. Correctif",
    "OT préparation <1 mois": "Prép. <1 mois",
    "OT préparation 1mois< <3mois": "Prép. 1-3 mois",
    "OT préparation >3 mois": "Prép. >3 mois",
    "OT planification <1 mois": "Planif. <1 mois",
    "OT planification 1mois< <3mois": "Planif. 1-3 mois",
    "OT planification >3 mois": "Planif. >3 mois",
    "OT exécution <1 mois": "Exéc. <1 mois",
    "OT exécution 1mois< <3mois": "Exéc. 1-3 mois",
    "OT exécution >3 mois": "Exéc. >3 mois",
    "Performance Graissage": "Perf. Graissage",
    "Performance Inspection": "Perf. Inspection",
    "Performance Systématiques": "Perf. Systématiques",
    "Taux d'approbation des Avis": "Approb. Avis",
    "OT LANC ESTIME": "OT Lancé Estimé",
    "Backlog préparation caractérisé": "Backlog Prép. Caract.",
    "Backlog planification caractérisé": "Backlog Planif. Caract.",
    "OT CONFIME": "OT Confirmé",
    "OT_COR_EGAL": "OT Coût Égal",
    "OT Fiabilité": "OT Fiabilité",
    "Total Avis de Panne": "Avis de Panne",
}

APP_URL = "https://tableaubordmc.streamlit.app"


def _color_for(val: float, cible: float, lower: bool) -> RGBColor:
    # CORRIGÉ : arrondi à l'entier AVANT comparaison (même précision que
    # l'affichage "f'{v:.0f}%'" dans les tableaux). Sans ça, une valeur
    # comme 99.52% s'affiche "100%" mais était comparée à la cible sans
    # arrondi (99.52 < 100) -> colorée orange alors que l'utilisateur voit
    # "100%" et s'attend à du vert. Avec l'arrondi, round(99.52)=100
    # correspond bien à la cible -> vert, cohérent avec l'affichage.
    val = round(val)
    if lower:
        if val <= cible:
            return GREEN
        if val <= cible * 1.5:
            return ORANGE
        return RED
    else:
        if val >= cible:
            return GREEN
        if val >= cible * 0.9:
            return ORANGE
        return RED


def _set_cell(cell, text, size=8, color=DARK, bold=False, bg=None, align=PP_ALIGN.LEFT):
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    if not p.runs:
        p.add_run()
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Arial"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Emu(18288)
    cell.margin_right = Emu(18288)
    cell.margin_top = Emu(0)
    cell.margin_bottom = Emu(0)
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE


def _remove_default_shadow(shape):
    """Desactive proprement l'effet d'ombre par defaut applique par
    PowerPoint/LibreOffice aux formes, SANS toucher au reste de spPr
    (contrairement a shape.shadow._element, qui pointe en fait sur le
    spPr entier quand aucune ombre n'existe deja -> le supprimer effacait
    aussi le remplissage/bordure qu'on venait de definir juste avant)."""
    spPr = shape._element.spPr
    # Ajoute un effectLst vide -> "pas d'effet" explicite, sans supprimer spPr.
    existing = spPr.find(qn('a:effectLst'))
    if existing is not None:
        spPr.remove(existing)
    from lxml import etree
    effect_lst = etree.SubElement(spPr, qn('a:effectLst'))


def build_poste_report_pptx(
    poste: str, pscore: float, qscore: float,
    kpi_perf: dict, kpi_qual: dict, cibles: dict,
    anomalies: dict, total_anomalies: int, plan_action: list,
    date_str: str,
) -> Presentation:
    """Construit et retourne l'objet Presentation (1 slide) pour un poste."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # layout vide

    def add_rect(x, y, w, h, color, shape_type=MSO_SHAPE.RECTANGLE, radius=None):
        sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        _remove_default_shadow(sh)
        if radius is not None and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
            try:
                sh.adjustments[0] = radius
            except Exception:
                pass
        return sh

    def add_text(x, y, w, h, text, size=10, color=DARK, bold=False, italic=False, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Arial"
        return tb

    # ══════════════════════════════════════════════════════════════
    # EN-TÊTE
    # ══════════════════════════════════════════════════════════════
    add_rect(0, 0, 13.333, 0.95, NAVY)
    add_text(0.35, 0.12, 5.5, 0.5, poste, size=26, color=WHITE, bold=True)
    add_text(0.35, 0.6, 6.5, 0.3, f"Rapport KPI Performance & Qualité — SAP PM OCP  •  {date_str}",
             size=10.5, color=RGBColor(0xCB, 0xD5, 0xE1))
    add_text(0.35, 0.99, 6.5, 0.16, f"Généré automatiquement depuis l'application {APP_URL}",
             size=7.5, color=RGBColor(0x94, 0xA3, 0xB8), italic=True)

    def badge(x, label, value, color):
        add_rect(x, 0.13, 2.05, 0.7, WHITE, MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
        add_text(x + 0.08, 0.17, 1.9, 0.24, label, size=8.5, bold=True, color=GREY)
        add_text(x + 0.08, 0.38, 1.9, 0.4, value, size=20, bold=True, color=color)

    p_color = GREEN if pscore >= 90 else (ORANGE if pscore >= 80 else RED)
    q_color = GREEN if qscore >= 90 else (ORANGE if qscore >= 80 else RED)
    badge(6.95, "SCORE PERFORMANCE", f"{pscore:.1f}%", p_color)
    badge(9.1, "SCORE QUALITÉ", f"{qscore:.1f}%", q_color)
    badge(11.25, "ANOMALIES", str(total_anomalies), RED)

    # ══════════════════════════════════════════════════════════════
    # Tableaux KPI (Performance / Qualité) — sans colonne Anomalies
    # ══════════════════════════════════════════════════════════════
    def kpi_table(x, y, w, title, kpi_dict, accent):
        add_text(x, y, w, 0.24, title, size=11, bold=True, color=accent)
        n = len(kpi_dict)
        rows, cols = n + 1, 3
        row_h = 0.185
        gts = slide.shapes.add_table(rows, cols, Inches(x), Inches(y + 0.26), Inches(w), Inches(row_h * rows))
        tbl = gts.table
        tbl.columns[0].width = Inches(w * 0.68)
        tbl.columns[1].width = Inches(w * 0.16)
        tbl.columns[2].width = Inches(w * 0.16)
        for j, h in enumerate(["Indicateur", "Val.", "Cible"]):
            _set_cell(tbl.cell(0, j), h, size=7.5, color=WHITE, bold=True, bg=accent,
                      align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
        for i, (k, v) in enumerate(kpi_dict.items(), start=1):
            cible = cibles[k]
            lower = k in LOWER_BETTER
            c = _color_for(v, cible, lower)
            _set_cell(tbl.cell(i, 0), k, size=7.3, color=DARK)
            _set_cell(tbl.cell(i, 1), f"{v:.0f}%", size=7.3, color=WHITE, bold=True, bg=c, align=PP_ALIGN.CENTER)
            _set_cell(tbl.cell(i, 2), f"{'≤' if lower else '≥'}{cible:.0f}", size=7, color=GREY, align=PP_ALIGN.CENTER)
        for r in tbl.rows:
            r.height = Inches(row_h)
        return tbl

    kpi_table(0.35, 1.12, 4.15, "INDICATEURS DE PERFORMANCE", kpi_perf, GREEN_ACCENT)
    kpi_table(4.7, 1.12, 4.15, "INDICATEURS DE QUALITÉ", kpi_qual, BLUE_ACCENT)

    # ══════════════════════════════════════════════════════════════
    # Graphique anomalies par indicateur (barres horizontales bleues)
    # ══════════════════════════════════════════════════════════════
    ano_entries = sorted([(k, v) for k, v in anomalies.items() if v > 0], key=lambda x: x[1])
    add_text(9.15, 1.12, 3.85, 0.24, "ANOMALIES PAR INDICATEUR", size=11, bold=True, color=BLUE_ACCENT)

    if ano_entries:
        chart_data = CategoryChartData()
        chart_data.categories = [SHORT_LABELS.get(k, k) for k, _ in ano_entries]
        chart_data.add_series("Anomalies", [v for _, v in ano_entries])
        gframe = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(9.15), Inches(1.38), Inches(3.85), Inches(3.55), chart_data
        )
        chart = gframe.chart
        chart.has_legend = False
        chart.has_title = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(6.8)
        dl.font.color.rgb = DARK
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
        series = plot.series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = BLUE_ACCENT
        plot.gap_width = 25
        cat_ax = chart.category_axis
        cat_ax.tick_labels.font.size = Pt(6.8)
        cat_ax.tick_labels.font.color.rgb = DARK
        cat_ax.has_major_gridlines = False
        cat_ax.format.line.fill.background()
        val_ax = chart.value_axis
        val_ax.visible = False
        val_ax.has_major_gridlines = False
    else:
        add_text(9.15, 1.5, 3.85, 0.4, "Aucune anomalie sur ce poste.", size=10, color=GREY, italic=True)

    # ══════════════════════════════════════════════════════════════
    # Suivi d'évolution (compact — historique limité pour l'instant)
    # ══════════════════════════════════════════════════════════════
    add_text(9.15, 5.1, 3.85, 0.24, "SUIVI D'ÉVOLUTION", size=11, bold=True, color=NAVY)
    add_rect(9.15, 5.35, 3.85, 1.85, LGREY, MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)

    evo_box = slide.shapes.add_textbox(Inches(9.3), Inches(5.45), Inches(3.6), Inches(1.65))
    tf = evo_box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0

    def evo_line(text_parts, first=False):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        for txt, color, bold, size in text_parts:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Arial"
        return p

    evo_line([("●  ", GREEN, True, 10), (f"Performance : {pscore:.1f}%", DARK, True, 10)], first=True)
    evo_line([("●  ", BLUE, True, 10), (f"Qualité : {qscore:.1f}%", DARK, True, 10)])
    evo_line([("●  ", RED, True, 10), (f"Anomalies : {total_anomalies}", DARK, True, 10)])
    evo_line([("", DARK, False, 6)])
    evo_line([("1 seule date enregistrée — tendance disponible dès la 2e extraction.", GREY, False, 8)])
    for p in tf.paragraphs:
        p.space_after = Pt(4)

    # ══════════════════════════════════════════════════════════════
    # Plan d'action — TOUS les indicateurs en anomalie
    # ══════════════════════════════════════════════════════════════
    plan_sorted = sorted(plan_action, key=lambda p: -p["nb_anom"])
    add_text(0.35, 4.55, 8.5, 0.24, f"PLAN D'ACTION — TOUS LES INDICATEURS EN ANOMALIE ({len(plan_sorted)})",
             size=11, bold=True, color=NAVY)

    n = len(plan_sorted)
    rows, cols = n + 1, 5
    row_h = 0.185
    gts = slide.shapes.add_table(rows, cols, Inches(0.35), Inches(4.82), Inches(8.5), Inches(row_h * rows))
    tbl = gts.table
    col_widths = [2.0, 0.6, 0.85, 1.05, 4.0]
    for j, cw in enumerate(col_widths):
        tbl.columns[j].width = Inches(cw)
    headers = ["Indicateur", "Écart", "Nécessité action", "Responsable", "Action"]
    for j, h in enumerate(headers):
        _set_cell(tbl.cell(0, j), h, size=6.8, color=WHITE, bold=True, bg=NAVY,
                  align=PP_ALIGN.LEFT if j in (0, 4) else PP_ALIGN.CENTER)
    for i, p_row in enumerate(plan_sorted, start=1):
        ecart = p_row["ecart"]
        needs_action = ecart < 0
        ecart_txt = f"{'+' if ecart > 0 else ''}{ecart}%"
        _set_cell(tbl.cell(i, 0), p_row["kpi"], size=6.8, color=DARK)
        _set_cell(tbl.cell(i, 1), ecart_txt, size=6.8, color=(RED if ecart < 0 else GREEN), bold=True, align=PP_ALIGN.CENTER)
        _set_cell(tbl.cell(i, 2), "Oui" if needs_action else "Non", size=7, color=WHITE, bold=True,
                  bg=(RED if needs_action else GREEN), align=PP_ALIGN.CENTER)
        _set_cell(tbl.cell(i, 3), p_row["responsable"], size=6.6, color=DARK)
        _set_cell(tbl.cell(i, 4), p_row["action"], size=6.0, color=GREY)
    for r in tbl.rows:
        r.height = Inches(row_h)

    return prs
